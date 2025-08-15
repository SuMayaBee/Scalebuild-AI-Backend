"""
Document Processing Service - Handles various file types for RAG system
"""
import os
import io
import tempfile
from typing import Dict, Any, Optional, BinaryIO
import PyPDF2
import docx
import pandas as pd
from pptx import Presentation
import json
import csv
from pathlib import Path

class DocumentProcessor:
    """Service to extract text from various document types"""
    
    SUPPORTED_FORMATS = {
        'pdf': 'PDF Document',
        'txt': 'Text File',
        'docx': 'Word Document',
        'doc': 'Word Document (Legacy)',
        'xlsx': 'Excel Spreadsheet',
        'xls': 'Excel Spreadsheet (Legacy)',
        'csv': 'CSV File',
        'pptx': 'PowerPoint Presentation',
        'ppt': 'PowerPoint Presentation (Legacy)',
        'json': 'JSON File',
        'md': 'Markdown File',
        'rtf': 'Rich Text Format'
    }
    
    def __init__(self):
        print("✅ Document processor initialized")
    
    def _clean_text(self, text: str) -> str:
        """
        Clean text content to remove problematic characters for database storage
        
        Args:
            text: Raw text content
            
        Returns:
            Cleaned text safe for database storage
        """
        if not text:
            return ""
        
        # Remove null characters (0x00) that cause PostgreSQL errors
        cleaned_text = text.replace('\x00', '')
        
        # Remove other problematic control characters but keep useful ones
        # Keep: \n (newline), \r (carriage return), \t (tab)
        # Remove: other control characters (0x01-0x08, 0x0B-0x0C, 0x0E-0x1F)
        import re
        cleaned_text = re.sub(r'[\x01-\x08\x0B-\x0C\x0E-\x1F\x7F]', '', cleaned_text)
        
        # Normalize whitespace - replace multiple consecutive whitespace with single space
        cleaned_text = re.sub(r'\s+', ' ', cleaned_text)
        
        # Remove excessive newlines (more than 2 consecutive)
        cleaned_text = re.sub(r'\n{3,}', '\n\n', cleaned_text)
        
        return cleaned_text.strip()
    
    def get_supported_formats(self) -> Dict[str, str]:
        """Get list of supported file formats"""
        return self.SUPPORTED_FORMATS
    
    def is_supported_format(self, filename: str) -> bool:
        """Check if file format is supported"""
        extension = Path(filename).suffix.lower().lstrip('.')
        return extension in self.SUPPORTED_FORMATS
    
    def extract_text_from_file(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """
        Extract text content from uploaded file
        
        Args:
            file_content: Binary content of the file
            filename: Original filename with extension
            
        Returns:
            Dictionary with extracted text and metadata
        """
        try:
            extension = Path(filename).suffix.lower().lstrip('.')
            
            if not self.is_supported_format(filename):
                raise ValueError(f"Unsupported file format: {extension}")
            
            print(f"📄 Processing {extension.upper()} file: {filename}")
            
            # Route to appropriate processor
            if extension == 'pdf':
                return self._process_pdf(file_content, filename)
            elif extension in ['txt', 'md']:
                return self._process_text(file_content, filename)
            elif extension == 'docx':
                return self._process_docx(file_content, filename)
            elif extension in ['xlsx', 'xls']:
                return self._process_excel(file_content, filename)
            elif extension == 'csv':
                return self._process_csv(file_content, filename)
            elif extension == 'pptx':
                return self._process_pptx(file_content, filename)
            elif extension == 'json':
                return self._process_json(file_content, filename)
            else:
                # Fallback to text processing
                return self._process_text(file_content, filename)
                
        except Exception as e:
            print(f"❌ Error processing file {filename}: {e}")
            raise
    
    def _process_pdf(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PDF file"""
        try:
            pdf_file = io.BytesIO(file_content)
            pdf_reader = PyPDF2.PdfReader(pdf_file)
            
            text_content = ""
            page_count = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages):
                try:
                    page_text = page.extract_text()
                    if page_text.strip():
                        text_content += f"\n--- Page {page_num + 1} ---\n"
                        text_content += page_text + "\n"
                except Exception as e:
                    print(f"⚠️ Error extracting page {page_num + 1}: {e}")
                    continue
            
            if not text_content.strip():
                raise ValueError("No text content could be extracted from PDF")
            
            metadata = {
                "file_type": "pdf",
                "page_count": page_count,
                "char_count": len(text_content),
                "extraction_method": "PyPDF2"
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Extracted {len(text_content)} characters from {page_count} pages")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters (removed {len(text_content) - len(cleaned_text)} problematic characters)")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing PDF: {e}")
            raise
    
    def _process_text(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Process plain text files"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'utf-16', 'latin-1', 'cp1252']
            text_content = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    text_content = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if text_content is None:
                raise ValueError("Could not decode text file with any supported encoding")
            
            metadata = {
                "file_type": "text",
                "encoding": used_encoding,
                "char_count": len(text_content),
                "line_count": len(text_content.splitlines())
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Processed text file with {len(text_content)} characters")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing text file: {e}")
            raise
    
    def _process_docx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from Word document"""
        try:
            doc_file = io.BytesIO(file_content)
            doc = docx.Document(doc_file)
            
            text_content = ""
            paragraph_count = 0
            
            for paragraph in doc.paragraphs:
                if paragraph.text.strip():
                    text_content += paragraph.text + "\n"
                    paragraph_count += 1
            
            # Extract text from tables
            table_count = 0
            for table in doc.tables:
                text_content += f"\n--- Table {table_count + 1} ---\n"
                for row in table.rows:
                    row_text = " | ".join([cell.text.strip() for cell in row.cells])
                    if row_text.strip():
                        text_content += row_text + "\n"
                table_count += 1
            
            if not text_content.strip():
                raise ValueError("No text content found in Word document")
            
            metadata = {
                "file_type": "docx",
                "paragraph_count": paragraph_count,
                "table_count": table_count,
                "char_count": len(text_content)
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Extracted text from Word document: {paragraph_count} paragraphs, {table_count} tables")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing Word document: {e}")
            raise
    
    def _process_excel(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from Excel spreadsheet"""
        try:
            excel_file = io.BytesIO(file_content)
            
            # Read all sheets
            excel_data = pd.read_excel(excel_file, sheet_name=None, engine='openpyxl')
            
            text_content = ""
            sheet_count = len(excel_data)
            total_rows = 0
            
            for sheet_name, df in excel_data.items():
                text_content += f"\n--- Sheet: {sheet_name} ---\n"
                
                # Convert DataFrame to text
                if not df.empty:
                    # Add column headers
                    headers = " | ".join(str(col) for col in df.columns)
                    text_content += f"Headers: {headers}\n"
                    
                    # Add data rows
                    for index, row in df.iterrows():
                        row_text = " | ".join(str(value) for value in row.values if pd.notna(value))
                        if row_text.strip():
                            text_content += f"Row {index + 1}: {row_text}\n"
                    
                    total_rows += len(df)
                else:
                    text_content += "Empty sheet\n"
            
            if not text_content.strip():
                raise ValueError("No data found in Excel file")
            
            metadata = {
                "file_type": "excel",
                "sheet_count": sheet_count,
                "total_rows": total_rows,
                "char_count": len(text_content)
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Processed Excel file: {sheet_count} sheets, {total_rows} total rows")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing Excel file: {e}")
            raise
    
    def _process_csv(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from CSV file"""
        try:
            # Try different encodings
            encodings = ['utf-8', 'latin-1', 'cp1252']
            csv_text = None
            used_encoding = None
            
            for encoding in encodings:
                try:
                    csv_text = file_content.decode(encoding)
                    used_encoding = encoding
                    break
                except UnicodeDecodeError:
                    continue
            
            if csv_text is None:
                raise ValueError("Could not decode CSV file")
            
            # Parse CSV
            csv_file = io.StringIO(csv_text)
            csv_reader = csv.reader(csv_file)
            
            text_content = ""
            row_count = 0
            headers = None
            
            for row_num, row in enumerate(csv_reader):
                if row_num == 0:
                    headers = row
                    text_content += f"Headers: {' | '.join(row)}\n"
                else:
                    row_text = " | ".join(str(cell) for cell in row if cell.strip())
                    if row_text.strip():
                        text_content += f"Row {row_num}: {row_text}\n"
                        row_count += 1
            
            if not text_content.strip():
                raise ValueError("No data found in CSV file")
            
            metadata = {
                "file_type": "csv",
                "encoding": used_encoding,
                "row_count": row_count,
                "column_count": len(headers) if headers else 0,
                "char_count": len(text_content)
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Processed CSV file: {row_count} rows, {len(headers) if headers else 0} columns")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing CSV file: {e}")
            raise
    
    def _process_pptx(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from PowerPoint presentation"""
        try:
            ppt_file = io.BytesIO(file_content)
            presentation = Presentation(ppt_file)
            
            text_content = ""
            slide_count = len(presentation.slides)
            
            for slide_num, slide in enumerate(presentation.slides):
                text_content += f"\n--- Slide {slide_num + 1} ---\n"
                
                # Extract text from shapes
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        text_content += shape.text + "\n"
                    
                    # Extract text from tables in slides
                    if shape.has_table:
                        table = shape.table
                        for row in table.rows:
                            row_text = " | ".join([cell.text.strip() for cell in row.cells])
                            if row_text.strip():
                                text_content += row_text + "\n"
            
            if not text_content.strip():
                raise ValueError("No text content found in PowerPoint presentation")
            
            metadata = {
                "file_type": "pptx",
                "slide_count": slide_count,
                "char_count": len(text_content)
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Extracted text from PowerPoint: {slide_count} slides")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing PowerPoint file: {e}")
            raise
    
    def _process_json(self, file_content: bytes, filename: str) -> Dict[str, Any]:
        """Extract text from JSON file"""
        try:
            json_text = file_content.decode('utf-8')
            json_data = json.loads(json_text)
            
            # Convert JSON to readable text
            text_content = self._json_to_text(json_data)
            
            metadata = {
                "file_type": "json",
                "char_count": len(text_content),
                "original_size": len(json_text)
            }
            
            # Clean text to remove null characters and other problematic content
            cleaned_text = self._clean_text(text_content)
            
            print(f"✅ Processed JSON file: {len(text_content)} characters")
            print(f"🧹 Cleaned text: {len(cleaned_text)} characters")
            
            return {
                "text": cleaned_text,
                "metadata": metadata
            }
            
        except Exception as e:
            print(f"❌ Error processing JSON file: {e}")
            raise
    
    def _json_to_text(self, data, prefix="", level=0) -> str:
        """Convert JSON data to readable text format"""
        text = ""
        indent = "  " * level
        
        if isinstance(data, dict):
            for key, value in data.items():
                current_prefix = f"{prefix}.{key}" if prefix else key
                
                if isinstance(value, (dict, list)):
                    text += f"{indent}{key}:\n"
                    text += self._json_to_text(value, current_prefix, level + 1)
                else:
                    text += f"{indent}{key}: {value}\n"
        
        elif isinstance(data, list):
            for i, item in enumerate(data):
                current_prefix = f"{prefix}[{i}]" if prefix else f"[{i}]"
                
                if isinstance(item, (dict, list)):
                    text += f"{indent}Item {i + 1}:\n"
                    text += self._json_to_text(item, current_prefix, level + 1)
                else:
                    text += f"{indent}Item {i + 1}: {item}\n"
        
        else:
            text += f"{indent}{data}\n"
        
        return text

# Global document processor instance
document_processor = DocumentProcessor()